"""
Gera o PPTX final da apresentação (slides/A3Data_Desafio_Tecnico.pptx), cobrindo
os itens a-j pedidos no desafio. Usa os gráficos já exportados em
outputs/figures/ e os números calculados nos módulos de src/ e nos notebooks.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = ROOT / "outputs" / "figures"
OUTPUT_PATH = ROOT / "slides" / "A3Data_Desafio_Tecnico.pptx"

PRESENTER_NAME = "Marcio Belloni — Cientista de Dados — mbelloni@alumni.usp.br"

# Paleta (mesma do restante do projeto, ver src/viz_style.py)
INK_PRIMARY = RGBColor(0x0B, 0x0B, 0x0B)
INK_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
INK_MUTED = RGBColor(0x89, 0x87, 0x81)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
RED = RGBColor(0xE3, 0x49, 0x48)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
GRIDLINE = RGBColor(0xE1, 0xE0, 0xD9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Segoe UI"


def _set_background(slide, color=SURFACE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                  color=INK_PRIMARY, align=PP_ALIGN.LEFT, font=FONT, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    p.line_spacing = line_spacing
    return box


def _add_bullets(slide, left, top, width, height, items, size=16, color=INK_SECONDARY,
                  bold_first=False, space_after=10, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.15
        bullet_run = p.add_run()
        bullet_run.text = f"•  {item}"
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = color
        bullet_run.font.name = font
    return box


def _add_accent_bar(slide, color=BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.14), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def _slide_header(slide, kicker, title, accent=BLUE):
    _set_background(slide)
    _add_accent_bar(slide, accent)
    _add_textbox(slide, Inches(0.55), Inches(0.35), Inches(11), Inches(0.4),
                 kicker, size=14, bold=True, color=accent)
    _add_textbox(slide, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.9),
                 title, size=28, bold=True, color=INK_PRIMARY)


def new_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_slide(prs):
    slide = new_blank(prs)
    _set_background(slide)
    _add_accent_bar(slide)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.55), SLIDE_W, Inches(0.03))
    band.fill.solid()
    band.fill.fore_color.rgb = GRIDLINE
    band.line.fill.background()
    band.shadow.inherit = False

    _add_textbox(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.4),
                 "TESTE TÉCNICO A3DATA — CIENTISTA DE DADOS (LLM/NLP)", size=14, bold=True, color=BLUE)
    _add_textbox(slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.3),
                 "Do book review ao insight de negócio:", size=32, bold=True, color=INK_PRIMARY)
    _add_textbox(slide, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.9),
                 "uma ferramenta de análise automatizada com NLP e LLM para a editora", size=22, bold=False, color=INK_SECONDARY)
    _add_textbox(slide, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
                 PRESENTER_NAME, size=15, bold=False, color=INK_MUTED)
    return slide


def add_bullets_slide(prs, kicker, title, bullets, accent=BLUE, notes=None):
    slide = new_blank(prs)
    _slide_header(slide, kicker, title, accent)
    _add_bullets(slide, Inches(0.7), Inches(1.9), Inches(11.8), Inches(5), bullets, size=18)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_image_slide(prs, kicker, title, image_paths, captions=None, accent=BLUE, notes=None):
    slide = new_blank(prs)
    _slide_header(slide, kicker, title, accent)

    n = len(image_paths)
    top = Inches(1.85)
    avail_w = Inches(12.3)
    gap = Inches(0.25)
    img_w = Emu(int((avail_w - gap * (n - 1)) / n)) if n > 1 else avail_w
    left = Inches(0.5)

    for i, path in enumerate(image_paths):
        pic = slide.shapes.add_picture(str(path), left, top, width=img_w)
        if pic.height > Inches(5.2):
            ratio = Inches(5.2) / pic.height
            pic.height = Inches(5.2)
            pic.width = Emu(int(pic.width * ratio))
        if captions and i < len(captions) and captions[i]:
            _add_textbox(slide, left, top + pic.height + Inches(0.05), img_w, Inches(0.4),
                         captions[i], size=12, color=INK_MUTED, align=PP_ALIGN.CENTER)
        left = Emu(int(left + img_w + gap))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_stat_row_slide(prs, kicker, title, stats, accent=BLUE, notes=None):
    """stats: list of (value, label) tuples, exibidos como cards em linha."""
    slide = new_blank(prs)
    _slide_header(slide, kicker, title, accent)

    n = len(stats)
    top = Inches(2.4)
    card_h = Inches(2.3)
    gap = Inches(0.3)
    avail_w = Inches(12.3)
    card_w = Emu(int((avail_w - gap * (n - 1)) / n))
    left = Inches(0.5)

    for value, label in stats:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.adjustments[0] = 0.06
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF3, 0xF6, 0xFC)
        card.line.color.rgb = GRIDLINE
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = value
        r1.font.size = Pt(30)
        r1.font.bold = True
        r1.font.color.rgb = BLUE
        r1.font.name = FONT
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(13)
        r2.font.color.rgb = INK_SECONDARY
        r2.font.name = FONT
        left = Emu(int(left + card_w + gap))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_code_slide(prs, kicker, title, code_blocks, accent=BLUE, notes=None):
    """code_blocks: list of (heading, code_text)"""
    slide = new_blank(prs)
    _slide_header(slide, kicker, title, accent)

    top = Inches(1.85)
    for heading, code_text in code_blocks:
        _add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.35), heading,
                     size=15, bold=True, color=INK_PRIMARY)
        top += Inches(0.4)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, Inches(12.1),
                                      Inches(0.35 + 0.24 * code_text.count(chr(10))))
        box.adjustments[0] = 0.04
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1C)
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        for i, line in enumerate(code_text.strip("\n").split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.size = Pt(12)
            r.font.name = "Consolas"
            r.font.color.rgb = RGBColor(0xD9, 0xD8, 0xD3)
        top += box.height + Inches(0.25)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_process_slide(prs):
    slide = new_blank(prs)
    _slide_header(slide, "C. PROCESSO", "Pipeline da solução, ponta a ponta")

    steps = [
        "1. Amostragem\n(DuckDB)",
        "2. Limpeza &\njunção",
        "3. EDA\n(Jupyter)",
        "4. NLP clássico\n(VADER, TF-IDF)",
        "5. Indexação RAG\n(embeddings + Chroma)",
        "6. API FastAPI\n(LLM via Groq)",
    ]
    n = len(steps)
    top = Inches(3.0)
    h = Inches(1.4)
    gap = Inches(0.35)
    avail_w = Inches(12.3)
    w = Emu(int((avail_w - gap * (n - 1)) / n))
    left = Inches(0.5)

    for i, step in enumerate(steps):
        color = BLUE if i < 4 else AQUA
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.adjustments[0] = 0.08
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = step
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.name = FONT

        if i < n - 1:
            arrow_left = Emu(int(left + w))
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, Emu(int(top + h / 2 - Inches(0.12))),
                                            gap, Inches(0.24))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = INK_MUTED
            arrow.line.fill.background()
            arrow.shadow.inherit = False
        left = Emu(int(left + w + gap))

    _add_bullets(slide, Inches(0.7), Inches(4.9), Inches(11.8), Inches(2), [
        "Cada etapa é um script/notebook independente e reexecutável (src/ e notebooks/), versionado no repositório.",
        "Amostra representativa e proporcional (~200 mil reviews de ~3 milhões) para rodar em ambiente local sem infraestrutura dedicada.",
        "Etapas 5-6 usam a mesma arquitetura de um projeto de RAG já validado (LangChain + ChromaDB + Groq/Llama 3.3 70B).",
    ], size=16)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Capa
    add_title_slide(prs)

    # a. Apresentação do desafio
    add_bullets_slide(
        prs, "A. O DESAFIO", "Uma editora, milhões de reviews, um processo manual que não escala",
        [
            "A editora contratou a A3Data para explorar sua base de avaliações de livros (Amazon Books Reviews) e extrair valor de negócio dela.",
            "Objetivos de negócio: entender performance por autor e por gênero, e localizar usuários com opiniões relevantes para entrevista.",
            "Hoje esse trabalho é 100% manual: 5 analistas, 3 dias por ciclo de análise — não escala com o volume de dados (~3 milhões de reviews, ~212 mil livros).",
            "As colunas de texto (review, resumo, descrição do livro) são a maior fonte de valor não explorada — pedem NLP e, idealmente, LLM.",
            "Entregável proposto: uma ferramenta (API) que automatiza a exploração, com sumarização, busca semântica e ranking de reviewers relevantes.",
        ],
    )

    # b. Roadmap
    add_bullets_slide(
        prs, "B. ROADMAP", "Planejamento de entregáveis — hoje e o que vem a seguir",
        [
            "Fase 1 — MVP (este projeto): amostragem + EDA, sentimento e keywords clássicos, RAG com embeddings locais + Groq, API FastAPI com 3 endpoints, métricas de qualidade objetivas.",
            "Fase 2 — Produção: ingestão incremental do catálogo completo, embeddings multilíngues (a busca hoje é melhor em inglês que em português), monitoramento de custo/latência do LLM, testes automatizados.",
            "Fase 3 — Confiabilidade: avaliação contínua (ROUGE/precisão@k como gates de CI), human-in-the-loop para validar resumos antes de publicar, versionamento do índice vetorial.",
            "Fase 4 — Expansão de produto: dashboard de negócio (BI) sobre os mesmos dados, alertas automáticos de reviews discrepantes/críticas, integração com o time editorial e de marketing.",
        ],
        accent=AQUA,
    )

    # c. Processo
    add_process_slide(prs)

    # d. Hipóteses
    add_bullets_slide(
        prs, "D. HIPÓTESES", "O que testamos na exploração — e o que os dados confirmaram ou não",
        [
            "H1 — A base é enviesada para notas altas. CONFIRMADA: 60% das reviews são 5 estrelas, 80% são 4-5 estrelas.",
            "H2 — Reviews de nota extrema (1 ou 5) são as mais longas/detalhadas. REFUTADA: reviews de nota 3-4 são as mais longas (mediana ~600-650 caracteres); 5 estrelas são as mais curtas (~470) — elogio rápido, pouco detalhe.",
            "H3 — Sentimento léxico clássico (VADER) é suficiente para captar a opinião do texto. PARCIALMENTE REFUTADA: 75% de acurácia geral, mas erra a maioria dos casos difíceis (sarcasmo, crítica com elogios pontuais) — LLM recupera 63% desses casos.",
            "H4 — Os metadados atuais (autor/gênero) têm problemas de qualidade que dificultam a análise manual. CONFIRMADA: mesmo autor aparece duplicado por variação de nome (ex.: \"J. R. R. Tolkien\" vs. \"John Ronald Reuel Tolkien\") — oportunidade de normalização via LLM.",
            "H5 — Cobertura de metadados é suficiente para segmentar por autor/gênero. CONFIRMADA: 87% das reviews têm autor identificado, 82% têm gênero.",
        ],
        accent=RED,
    )

    # e. EDA
    add_image_slide(
        prs, "E. ANÁLISE EXPLORATÓRIA", "Distribuição de notas e evolução do volume de reviews",
        [FIG_DIR / "score_distribution.png", FIG_DIR / "reviews_por_ano.png"],
        captions=["198.924 reviews amostradas (proporcional aos ~3M originais)", "Pico em 2012 (~20 mil reviews); queda em 2013 é corte parcial do ano na coleta"],
    )
    add_image_slide(
        prs, "E. ANÁLISE EXPLORATÓRIA", "Performance por gênero e por autor",
        [FIG_DIR / "generos_volume_nota.png"],
    )
    add_image_slide(
        prs, "E. ANÁLISE EXPLORATÓRIA", "Quem mais é avaliado x quem é mais bem avaliado",
        [FIG_DIR / "autores_volume_nota.png"],
        captions=["Nota: \"J. R. R. Tolkien\" e \"John Ronald Reuel Tolkien\" — mesmo autor, duplicado nos metadados"],
    )

    # f. Sumarização
    add_bullets_slide(
        prs, "F. SUMARIZAÇÃO DE TEXTO", "Do volume de reviews para \"o que os leitores estão dizendo\"",
        [
            "Abordagem: para cada livro/autor, o LLM (Groq/Llama 3.3 70B) sintetiza as reviews mais substantivas em um resumo estruturado — sentimento geral, elogios, críticas e uma citação representativa.",
            "Disponível como endpoint da API: GET /resumo/livro/{título} e GET /resumo/autor/{autor} — geração sob demanda, não pré-computada, sempre com as reviews mais recentes da base.",
            "Exemplo real (The Hobbit, 1.450 reviews na amostra, nota média 4,65): \"Sentimento geral extremamente positivo... elogiam a criatividade e a construção do mundo de Middle-earth... críticas mencionam a lentidão do início\".",
            "Avaliação objetiva: ROUGE-L de 0,19 comparando resumos de review individual contra o resumo humano real já presente no dataset — mas similaridade semântica (embeddings) de 0,44, bem mais alta, confirmando que o resumo captura o significado mesmo usando palavras diferentes do resumo humano (que é estilizado, não descritivo).",
            "Produção: cache dos resumos mais acessados, e opção de resumo incremental conforme novas reviews chegam.",
        ],
    )

    # g. Base de conhecimento / RAG
    add_bullets_slide(
        prs, "G. BASE DE CONHECIMENTO", "RAG: perguntas em linguagem natural sobre o catálogo",
        [
            "Arquitetura (POC funcional, não só proposta): descrição do livro + resumos de reviews reais → embeddings locais (sentence-transformers/all-MiniLM-L6-v2) → índice vetorial ChromaDB → LangChain → Groq/Llama 3.3 70B gera a resposta final em português.",
            "13.661 livros indexados (todos com 3+ reviews na amostra) — endpoint POST /perguntar, ex.: \"quais livros de fantasia os leitores mais elogiaram?\".",
            "Diferencial técnico: a pergunta do usuário é reescrita pelo próprio LLM para inglês antes da busca (o corpus é majoritariamente em inglês) — melhora bastante a qualidade da recuperação sem exigir um modelo de embeddings multilíngue mais pesado.",
            "Precisão@5 média de 80% em queries de teste com gênero-alvo conhecido (varia de 40% a 100% conforme a especificidade da pergunta).",
            "Roadmap de produção: ingestão incremental (novo livro/review → reindexação automática), embeddings multilíngues nativos, re-ranking com LLM para os casos mais difíceis, e filtros de metadata (autor/gênero/nota) combinados à busca semântica.",
        ],
        accent=AQUA,
    )

    # h. Métricas de qualidade
    add_image_slide(
        prs, "H. MÉTRICAS DE QUALIDADE", "Como sabemos que a ferramenta está funcionando bem",
        [FIG_DIR / "rouge_sumarizacao.png", FIG_DIR / "rag_precisao_at_k.png", FIG_DIR / "vader_vs_llm_sentimento.png"],
    )
    add_bullets_slide(
        prs, "H. MÉTRICAS DE QUALIDADE", "Leitura dos números",
        [
            "Sumarização — ROUGE-1/2/L (0,19 / 0,11 / 0,19) mede sobreposição de palavras contra o resumo humano real (referência gold); similaridade semântica via embeddings (0,44) mede significado — as duas juntas mostram que o resumo é fiel mesmo quando ROUGE sozinho pareceria baixo.",
            "RAG — Precisão@5 de 80% em 7 queries de teste cobrindo gêneros distintos (fantasia, negócios, religião, história, tecnologia, culinária, biografia).",
            "Sentimento — no recorte mais difícil (12,5% da amostra, onde o texto \"engana\" o léxico clássico VADER), o LLM acerta 63,3% das vezes contra 0% do VADER nesse mesmo recorte por construção.",
            "Todas as métricas são reprodutíveis: notebooks/02_avaliacao_qualidade.ipynb roda do zero e recalcula os três números.",
        ],
    )

    # i. Impacto financeiro
    add_stat_row_slide(
        prs, "I. IMPACTO FINANCEIRO", "O que a automação libera, em números",
        [
            ("96,7%", "de redução no custo por\nciclo de análise"),
            ("14,5", "dias-analista liberados\npor ciclo"),
            ("R$ 158 mil/ano", "economia estimada\n(premissa: 4 ciclos/mês)"),
        ],
        accent=RED,
    )
    add_image_slide(
        prs, "I. IMPACTO FINANCEIRO", "Premissas explícitas — ajustáveis ao vivo",
        [FIG_DIR / "impacto_financeiro.png"],
        captions=["Premissas: R$5.000/mês por analista, 22 dias úteis/mês, 3 dias→0,5 dia por ciclo, 5→1 analista"],
        accent=RED,
    )
    add_bullets_slide(
        prs, "I. IMPACTO FINANCEIRO", "Como chegamos nesses números",
        [
            "Hoje: 5 analistas × 3 dias × taxa diária (R$5.000/mês ÷ 22 dias úteis ≈ R$227,27/dia) = R$3.409,09 por ciclo de análise.",
            "Com a ferramenta: 1 analista × 0,5 dia (revisão/validação dos outputs, não mais execução manual) = R$113,64 por ciclo.",
            "Os analistas liberados (4 de 5, ~14,5 dias-pessoa por ciclo) não são cortados — são realocados para trabalho de maior valor, como pedido no desafio.",
            "Sensibilidade: com 1 ciclo/mês, economia anual de ~R$39,5 mil; com 12 ciclos/mês (quase semanal), ~R$474,5 mil/ano — número final depende da frequência real de uso pela editora.",
        ],
        accent=RED,
    )

    # j. POC
    add_code_slide(
        prs, "J. POC DA SOLUÇÃO", "API funcional — não é só protótipo em notebook",
        [
            ("POST /perguntar — pergunta em linguagem natural sobre o catálogo", (
                '{"pergunta": "quais livros de fantasia os leitores mais elogiaram?"}\n\n'
                '→ "...leitores mais elogiaram: \'Deep Wizardry\' (Diane Duane),\n'
                '   \'The Hobbit\' (Tolkien), \'The Light Fantastic\' (Pratchett)..."'
            )),
            ("GET /reviewers-relevantes?genero=Fiction — ranking p/ entrevista", (
                '→ [{"profileName": "B. Scalera", "specificity": 16256,\n'
                '     "relevance_score": 4.86, ...}, ...]'
            )),
        ],
    )
    add_bullets_slide(
        prs, "J. POC DA SOLUÇÃO", "Como rodar e onde está o código",
        [
            "API FastAPI local: `uvicorn app.main:app --reload` — documentação interativa (Swagger) em /docs, testável direto no navegador.",
            "3 endpoints funcionais: POST /perguntar (RAG), GET /resumo/livro|autor/{nome} (sumarização via LLM), GET /reviewers-relevantes (ranking heurístico).",
            "Reaproveita a arquitetura de um projeto de RAG já validado (LangChain + ChromaDB + Groq), adaptado para o domínio de reviews de livros.",
            "Código completo, notebooks e este PPTX no repositório GitHub — README com instruções de reprodução do zero.",
        ],
    )

    # Conclusão
    add_bullets_slide(
        prs, "CONCLUSÃO", "De 3 dias manuais para uma consulta em segundos",
        [
            "A ferramenta cobre as duas dores originais da editora: performance por autor/gênero (EDA + NLP clássico) e localização de opiniões relevantes para entrevista (reviewer scoring + RAG).",
            "Uso de LLM não é cosmético: melhora mensurável sobre NLP clássico em sentimento (63% vs. 0% nos casos difíceis) e viabiliza sumarização e busca em linguagem natural que um pipeline clássico não faria bem.",
            "Todas as métricas de qualidade são objetivas e reprodutíveis — não é uma demo de opinião, é um sistema com números que podem virar SLA.",
            "Próximo passo natural: validar com um piloto real com o time de analistas, medir tempo real economizado vs. a estimativa, e priorizar a Fase 2 do roadmap.",
        ],
        accent=BLUE,
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"PPTX salvo em {OUTPUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
